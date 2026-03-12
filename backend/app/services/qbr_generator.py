"""QBR (Quarterly Business Review) generation service."""

import io
import logging
from datetime import date, datetime, timedelta
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from fpdf import FPDF

logger = logging.getLogger(__name__)


class QBRGenerator:
    """Generate QBR documents in PowerPoint or PDF format."""
    
    # Brand colors
    PRIMARY_COLOR = RGBColor(0x0F, 0x76, 0x6E)  # Teal
    SECONDARY_COLOR = RGBColor(0x1E, 0x29, 0x3B)  # Dark slate
    ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)  # Emerald
    ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)  # Amber
    ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)  # Red
    
    def __init__(self, account_data: dict, support_data: dict, usage_data: dict, health_data: dict):
        """
        Initialize QBR generator with account data.
        
        Args:
            account_data: Account details (name, industry, csm, ae, arr, renewals)
            support_data: Freshdesk support analysis
            usage_data: Pendo usage analysis
            health_data: Health score breakdown
        """
        self.account = account_data
        self.support = support_data
        self.usage = usage_data
        self.health = health_data
        self.quarter = self._get_current_quarter()
    
    def _get_current_quarter(self) -> str:
        """Get current quarter string (e.g., 'Q1 2026')."""
        today = date.today()
        quarter = (today.month - 1) // 3 + 1
        return f"Q{quarter} {today.year}"
    
    def _get_health_color(self, score: int) -> RGBColor:
        """Get color based on health score."""
        if score >= 70:
            return self.ACCENT_GREEN
        elif score >= 40:
            return self.ACCENT_AMBER
        return self.ACCENT_RED
    
    def _get_health_label(self, score: int) -> str:
        """Get label based on health score."""
        if score >= 70:
            return "Good"
        elif score >= 40:
            return "At Risk"
        return "Critical"
    
    def generate_pptx(self) -> bytes:
        """Generate PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Add slides
        self._add_title_slide(prs)
        self._add_executive_summary_slide(prs)
        self._add_usage_overview_slide(prs)
        self._add_support_overview_slide(prs)
        self._add_health_breakdown_slide(prs)
        self._add_recommendations_slide(prs)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
    
    def _add_title_slide(self, prs: Presentation):
        """Add title slide."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.SECONDARY_COLOR
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Quarterly Business Review"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # Account name
        account_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12), Inches(0.8))
        tf = account_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.account.get("name", "Account")
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Quarter
        quarter_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.5))
        tf = quarter_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.quarter
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.CENTER
        
        # Generated date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Generated: {date.today().strftime('%B %d, %Y')}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_executive_summary_slide(self, prs: Presentation):
        """Add executive summary slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Title
        self._add_slide_title(slide, "Executive Summary")
        
        # Health Score Card
        health_score = self.health.get("score", 0)
        health_label = self._get_health_label(health_score)
        health_color = self._get_health_color(health_score)
        
        # Health score box
        health_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5), Inches(3), Inches(2)
        )
        health_box.fill.solid()
        health_box.fill.fore_color.rgb = health_color
        health_box.line.fill.background()
        
        health_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(3), Inches(1))
        tf = health_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(health_score)
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        health_label_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(3), Inches(0.5))
        tf = health_label_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Health: {health_label}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # Key Metrics
        metrics_y = 1.5
        metrics_x = 4.0
        
        metrics = [
            ("Industry", self.account.get("industry", "N/A")),
            ("CSM", self.account.get("csm_name", "Unassigned")),
            ("Account Executive", self.account.get("ae_name", "Unassigned")),
            ("Active Users (30d)", str(self.usage.get("pendo_summary", {}).get("current_active_visitors", 0))),
            ("Open Tickets", str(self.support.get("open_tickets", 0))),
            ("Critical Tickets", str(self.support.get("critical_tickets", 0))),
        ]
        
        for i, (label, value) in enumerate(metrics):
            y_pos = metrics_y + (i * 0.7)
            
            label_box = slide.shapes.add_textbox(Inches(metrics_x), Inches(y_pos), Inches(3), Inches(0.35))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            
            value_box = slide.shapes.add_textbox(Inches(metrics_x + 3.5), Inches(y_pos), Inches(3), Inches(0.35))
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_COLOR
        
        # Renewal info
        renewals = self.account.get("renewals", [])
        if renewals:
            renewal_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(0.5))
            tf = renewal_box.text_frame
            p = tf.paragraphs[0]
            nearest = min(renewals, key=lambda r: r.get("renewal_days", 9999))
            days = nearest.get("renewal_days", "N/A")
            p.text = f"Next Renewal: {days} days" if isinstance(days, int) else "No upcoming renewals"
            p.font.size = Pt(16)
            p.font.color.rgb = self.SECONDARY_COLOR
    
    def _add_usage_overview_slide(self, prs: Presentation):
        """Add Pendo usage overview slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        self._add_slide_title(slide, "Product Usage — Pendo")
        
        pendo_summary = self.usage.get("pendo_summary", {})
        has_pendo = self.usage.get("has_pendo_data", False)
        
        if not has_pendo:
            no_data_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1))
            tf = no_data_box.text_frame
            p = tf.paragraphs[0]
            p.text = "No Pendo integration for this account"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            p.alignment = PP_ALIGN.CENTER
            return
        
        # Usage metrics cards
        metrics = [
            ("Active Users", 
             int(pendo_summary.get("current_active_visitors", 0)),
             pendo_summary.get("visitors_change_pct", 0),
             "Last 30 days"),
            ("Time Spent", 
             f"{int(pendo_summary.get('current_minutes', 0))}m",
             pendo_summary.get("minutes_change_pct", 0),
             "Total minutes"),
            ("Events", 
             int(pendo_summary.get("current_events", 0)),
             pendo_summary.get("events_change_pct", 0),
             "Interactions"),
        ]
        
        for i, (label, value, change, subtitle) in enumerate(metrics):
            x_pos = 0.5 + (i * 4.2)
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(1.5), Inches(3.8), Inches(2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            
            # Label
            label_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.7), Inches(3.4), Inches(0.4))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            
            # Value
            value_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.1), Inches(3.4), Inches(0.7))
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_COLOR
            
            # Change indicator
            change_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.8), Inches(3.4), Inches(0.4))
            tf = change_box.text_frame
            p = tf.paragraphs[0]
            arrow = "↑" if change > 0 else "↓" if change < 0 else "→"
            p.text = f"{arrow} {abs(change):.1f}% vs prior period"
            p.font.size = Pt(11)
            p.font.color.rgb = self.ACCENT_GREEN if change >= 0 else self.ACCENT_RED
        
        # Top features section
        features = self.usage.get("pendo_features", {}).get("top_items", [])[:5]
        if features:
            features_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(6), Inches(0.5))
            tf = features_title.text_frame
            p = tf.paragraphs[0]
            p.text = "Top Features Used"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_COLOR
            
            for i, feature in enumerate(features):
                feat_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5 + i * 0.4), Inches(8), Inches(0.35))
                tf = feat_box.text_frame
                p = tf.paragraphs[0]
                name = feature.get("feature_name", feature.get("name", "Unknown"))
                clicks = feature.get("count_clicks", feature.get("clicks", 0))
                p.text = f"• {name} ({clicks} clicks)"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
    
    def _add_support_overview_slide(self, prs: Presentation):
        """Add Freshdesk support overview slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        self._add_slide_title(slide, "Support Analysis — Freshdesk")
        
        has_freshdesk = self.support.get("total_tickets", 0) > 0 or self.support.get("open_tickets", 0) > 0
        
        if not has_freshdesk:
            no_data_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1))
            tf = no_data_box.text_frame
            p = tf.paragraphs[0]
            p.text = "No Freshdesk data for this account"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            p.alignment = PP_ALIGN.CENTER
            return
        
        # Ticket metrics
        metrics = [
            ("Open Tickets", self.support.get("open_tickets", 0), RGBColor(0x3B, 0x82, 0xF6)),
            ("Critical", self.support.get("critical_tickets", 0), self.ACCENT_RED),
            ("High Priority", self.support.get("high_tickets", 0), self.ACCENT_AMBER),
            ("Total (All Time)", self.support.get("total_tickets", 0), RGBColor(0x64, 0x74, 0x8B)),
        ]
        
        for i, (label, value, color) in enumerate(metrics):
            x_pos = 0.5 + (i * 3.2)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(1.5), Inches(2.9), Inches(1.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            
            value_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.7), Inches(2.5), Inches(0.7))
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = color
            
            label_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.4), Inches(2.5), Inches(0.4))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        
        # Resolution stats
        resolution_stats = self.support.get("resolution_stats", {})
        if resolution_stats:
            res_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(6), Inches(0.5))
            tf = res_title.text_frame
            p = tf.paragraphs[0]
            p.text = "Resolution Time"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_COLOR
            
            median_days = resolution_stats.get("median_days", 0)
            mean_days = resolution_stats.get("mean_days", 0)
            
            res_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(6), Inches(0.8))
            tf = res_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Median: {median_days:.1f} days  |  Mean: {mean_days:.1f} days"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
        
        # Sentiment
        sentiment_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(6), Inches(0.5))
        tf = sentiment_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Ticket Sentiment"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.SECONDARY_COLOR
        
        positive = self.support.get("positive_ticket_count", 0)
        neutral = self.support.get("neutral_ticket_count", 0)
        negative = self.support.get("negative_ticket_count", 0)
        
        sentiment_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(10), Inches(0.5))
        tf = sentiment_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Positive: {positive}  |  Neutral: {neutral}  |  Negative: {negative}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
    
    def _add_health_breakdown_slide(self, prs: Presentation):
        """Add health score breakdown slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        self._add_slide_title(slide, "Health Score Breakdown")
        
        health_score = self.health.get("score", 0)
        factors = self.health.get("factors", [])
        
        # Main score display
        score_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.5), Inches(1.5), Inches(2), Inches(2)
        )
        score_box.fill.solid()
        score_box.fill.fore_color.rgb = self._get_health_color(health_score)
        score_box.line.fill.background()
        
        score_text = slide.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(2), Inches(1))
        tf = score_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(health_score)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # Factors breakdown
        factors_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(6), Inches(0.5))
        tf = factors_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Contributing Factors"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.SECONDARY_COLOR
        
        for i, factor in enumerate(factors[:4]):  # Max 4 factors
            y_pos = 4.6 + (i * 0.6)
            
            name = factor.get("name", "Unknown")
            points = factor.get("points", 0)
            detail = factor.get("detail", "")
            icon = factor.get("icon", "")
            
            factor_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(8), Inches(0.5))
            tf = factor_box.text_frame
            p = tf.paragraphs[0]
            deduction_text = f"-{points}" if points > 0 else "✓"
            p.text = f"{icon} {name}: {detail} ({deduction_text})"
            p.font.size = Pt(13)
            p.font.color.rgb = self.ACCENT_RED if points > 10 else self.ACCENT_AMBER if points > 0 else self.ACCENT_GREEN
        
        # Data sources
        has_pendo = self.health.get("has_pendo", False)
        has_freshdesk = self.health.get("has_freshdesk", False)
        
        sources_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(10), Inches(0.4))
        tf = sources_box.text_frame
        p = tf.paragraphs[0]
        pendo_status = "✓ Connected" if has_pendo else "✗ Not connected"
        freshdesk_status = "✓ Connected" if has_freshdesk else "✗ Not connected"
        p.text = f"Data Sources:  Pendo: {pendo_status}  |  Freshdesk: {freshdesk_status}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    def _add_recommendations_slide(self, prs: Presentation):
        """Add recommendations and next steps slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        self._add_slide_title(slide, "Recommendations & Next Steps")
        
        # Generate recommendations based on health factors
        recommendations = []
        
        health_score = self.health.get("score", 100)
        factors = self.health.get("factors", [])
        
        for factor in factors:
            points = factor.get("points", 0)
            name = factor.get("name", "")
            
            if points >= 20 and "Renewal" in name:
                recommendations.append("🔴 Critical: Schedule renewal discussion immediately")
            elif points >= 10 and "Renewal" in name:
                recommendations.append("🟡 Priority: Begin renewal planning process")
            
            if points >= 25 and "Pendo" in name:
                recommendations.append("🔴 Critical: Investigate product usage decline - schedule check-in call")
            elif points >= 10 and "Pendo" in name:
                recommendations.append("🟡 Monitor: Track usage trends over next 30 days")
            
            if points >= 15 and "Freshdesk" in name:
                recommendations.append("🔴 Action: Address critical support tickets")
            elif points >= 5 and "Freshdesk" in name:
                recommendations.append("🟡 Review: Follow up on open support cases")
        
        if health_score >= 70:
            recommendations.append("🟢 Opportunity: Explore expansion/upsell opportunities")
            recommendations.append("🟢 Engage: Consider case study or reference program")
        
        if not recommendations:
            recommendations = [
                "✓ Account health is stable",
                "• Continue regular engagement cadence",
                "• Monitor for any changes in usage or support needs",
            ]
        
        for i, rec in enumerate(recommendations[:8]):  # Max 8 recommendations
            y_pos = 1.5 + (i * 0.6)
            
            rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.5))
            tf = rec_box.text_frame
            p = tf.paragraphs[0]
            p.text = rec
            p.font.size = Pt(16)
            p.font.color.rgb = self.SECONDARY_COLOR
        
        # Contact info footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        csm = self.account.get("csm_name", "CSM")
        p.text = f"Questions? Contact your CSM: {csm}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_slide_title(self, slide, title: str):
        """Add consistent title to a slide."""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.SECONDARY_COLOR
        
        # Underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2), Inches(2), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.PRIMARY_COLOR
        line.line.fill.background()
    
    def generate_pdf(self) -> bytes:
        """Generate PDF document."""
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add pages
        self._add_pdf_title_page(pdf)
        self._add_pdf_executive_summary(pdf)
        self._add_pdf_usage_page(pdf)
        self._add_pdf_support_page(pdf)
        self._add_pdf_health_page(pdf)
        self._add_pdf_recommendations_page(pdf)
        
        return pdf.output()
    
    def _add_pdf_title_page(self, pdf: FPDF):
        """Add PDF title page."""
        pdf.add_page()
        pdf.set_fill_color(30, 41, 59)
        pdf.rect(0, 0, 210, 297, 'F')
        
        pdf.set_text_color(255, 255, 255)
        pdf.set_font('Helvetica', 'B', 36)
        pdf.set_y(100)
        pdf.cell(0, 20, 'Quarterly Business Review', align='C', ln=True)
        
        pdf.set_font('Helvetica', '', 24)
        pdf.set_text_color(15, 118, 110)
        pdf.cell(0, 15, self.account.get("name", "Account"), align='C', ln=True)
        
        pdf.set_font('Helvetica', '', 18)
        pdf.set_text_color(148, 163, 184)
        pdf.cell(0, 10, self.quarter, align='C', ln=True)
        
        pdf.set_y(260)
        pdf.set_font('Helvetica', '', 12)
        pdf.set_text_color(100, 116, 139)
        pdf.cell(0, 10, f"Generated: {date.today().strftime('%B %d, %Y')}", align='C')
    
    def _add_pdf_executive_summary(self, pdf: FPDF):
        """Add PDF executive summary page."""
        pdf.add_page()
        pdf.set_text_color(30, 41, 59)
        
        pdf.set_font('Helvetica', 'B', 20)
        pdf.cell(0, 15, 'Executive Summary', ln=True)
        pdf.set_draw_color(15, 118, 110)
        pdf.line(10, pdf.get_y(), 60, pdf.get_y())
        pdf.ln(10)
        
        # Health score
        health_score = self.health.get("score", 0)
        health_label = self._get_health_label(health_score)
        
        pdf.set_font('Helvetica', 'B', 14)
        pdf.cell(0, 10, f'Health Score: {health_score} ({health_label})', ln=True)
        pdf.ln(5)
        
        # Key metrics
        pdf.set_font('Helvetica', '', 12)
        metrics = [
            ("Industry", self.account.get("industry", "N/A")),
            ("CSM", self.account.get("csm_name", "Unassigned")),
            ("Account Executive", self.account.get("ae_name", "Unassigned")),
            ("Active Users (30d)", str(self.usage.get("pendo_summary", {}).get("current_active_visitors", 0))),
            ("Open Tickets", str(self.support.get("open_tickets", 0))),
        ]
        
        for label, value in metrics:
            pdf.cell(60, 8, f'{label}:', ln=False)
            pdf.set_font('Helvetica', 'B', 12)
            pdf.cell(0, 8, str(value), ln=True)
            pdf.set_font('Helvetica', '', 12)
    
    def _add_pdf_usage_page(self, pdf: FPDF):
        """Add PDF usage page."""
        pdf.add_page()
        pdf.set_text_color(30, 41, 59)
        
        pdf.set_font('Helvetica', 'B', 20)
        pdf.cell(0, 15, 'Product Usage - Pendo', ln=True)
        pdf.set_draw_color(15, 118, 110)
        pdf.line(10, pdf.get_y(), 60, pdf.get_y())
        pdf.ln(10)
        
        pendo_summary = self.usage.get("pendo_summary", {})
        has_pendo = self.usage.get("has_pendo_data", False)
        
        if not has_pendo:
            pdf.set_font('Helvetica', 'I', 14)
            pdf.set_text_color(148, 163, 184)
            pdf.cell(0, 20, 'No Pendo integration for this account', align='C')
            return
        
        pdf.set_font('Helvetica', '', 12)
        
        metrics = [
            ("Active Users (30d)", int(pendo_summary.get("current_active_visitors", 0)), pendo_summary.get("visitors_change_pct", 0)),
            ("Time Spent (minutes)", int(pendo_summary.get("current_minutes", 0)), pendo_summary.get("minutes_change_pct", 0)),
            ("Events", int(pendo_summary.get("current_events", 0)), pendo_summary.get("events_change_pct", 0)),
        ]
        
        for label, value, change in metrics:
            pdf.cell(70, 8, f'{label}:', ln=False)
            pdf.set_font('Helvetica', 'B', 12)
            arrow = "↑" if change > 0 else "↓" if change < 0 else "→"
            pdf.cell(0, 8, f'{value} ({arrow} {abs(change):.1f}%)', ln=True)
            pdf.set_font('Helvetica', '', 12)
    
    def _add_pdf_support_page(self, pdf: FPDF):
        """Add PDF support page."""
        pdf.add_page()
        pdf.set_text_color(30, 41, 59)
        
        pdf.set_font('Helvetica', 'B', 20)
        pdf.cell(0, 15, 'Support Analysis - Freshdesk', ln=True)
        pdf.set_draw_color(15, 118, 110)
        pdf.line(10, pdf.get_y(), 60, pdf.get_y())
        pdf.ln(10)
        
        pdf.set_font('Helvetica', '', 12)
        
        metrics = [
            ("Open Tickets", self.support.get("open_tickets", 0)),
            ("Critical Tickets", self.support.get("critical_tickets", 0)),
            ("High Priority Tickets", self.support.get("high_tickets", 0)),
            ("Total Tickets (All Time)", self.support.get("total_tickets", 0)),
        ]
        
        for label, value in metrics:
            pdf.cell(70, 8, f'{label}:', ln=False)
            pdf.set_font('Helvetica', 'B', 12)
            pdf.cell(0, 8, str(value), ln=True)
            pdf.set_font('Helvetica', '', 12)
        
        pdf.ln(5)
        
        # Resolution stats
        resolution_stats = self.support.get("resolution_stats", {})
        if resolution_stats:
            pdf.set_font('Helvetica', 'B', 14)
            pdf.cell(0, 10, 'Resolution Time', ln=True)
            pdf.set_font('Helvetica', '', 12)
            pdf.cell(0, 8, f"Median: {resolution_stats.get('median_days', 0):.1f} days  |  Mean: {resolution_stats.get('mean_days', 0):.1f} days", ln=True)
    
    def _add_pdf_health_page(self, pdf: FPDF):
        """Add PDF health breakdown page."""
        pdf.add_page()
        pdf.set_text_color(30, 41, 59)
        
        pdf.set_font('Helvetica', 'B', 20)
        pdf.cell(0, 15, 'Health Score Breakdown', ln=True)
        pdf.set_draw_color(15, 118, 110)
        pdf.line(10, pdf.get_y(), 60, pdf.get_y())
        pdf.ln(10)
        
        health_score = self.health.get("score", 0)
        health_label = self._get_health_label(health_score)
        factors = self.health.get("factors", [])
        
        pdf.set_font('Helvetica', 'B', 24)
        pdf.cell(0, 15, f'Score: {health_score} ({health_label})', ln=True)
        pdf.ln(5)
        
        pdf.set_font('Helvetica', 'B', 14)
        pdf.cell(0, 10, 'Contributing Factors:', ln=True)
        pdf.set_font('Helvetica', '', 12)
        
        for factor in factors:
            name = factor.get("name", "Unknown")
            points = factor.get("points", 0)
            detail = factor.get("detail", "")
            deduction = f"-{points}" if points > 0 else "✓"
            pdf.cell(0, 8, f"• {name}: {detail} ({deduction})", ln=True)
    
    def _add_pdf_recommendations_page(self, pdf: FPDF):
        """Add PDF recommendations page."""
        pdf.add_page()
        pdf.set_text_color(30, 41, 59)
        
        pdf.set_font('Helvetica', 'B', 20)
        pdf.cell(0, 15, 'Recommendations & Next Steps', ln=True)
        pdf.set_draw_color(15, 118, 110)
        pdf.line(10, pdf.get_y(), 60, pdf.get_y())
        pdf.ln(10)
        
        # Generate recommendations
        recommendations = []
        health_score = self.health.get("score", 100)
        factors = self.health.get("factors", [])
        
        for factor in factors:
            points = factor.get("points", 0)
            name = factor.get("name", "")
            
            if points >= 20 and "Renewal" in name:
                recommendations.append("Critical: Schedule renewal discussion immediately")
            elif points >= 10 and "Renewal" in name:
                recommendations.append("Priority: Begin renewal planning process")
            
            if points >= 25 and "Pendo" in name:
                recommendations.append("Critical: Investigate product usage decline")
            elif points >= 10 and "Pendo" in name:
                recommendations.append("Monitor: Track usage trends over next 30 days")
            
            if points >= 15 and "Freshdesk" in name:
                recommendations.append("Action: Address critical support tickets")
        
        if health_score >= 70:
            recommendations.append("Opportunity: Explore expansion/upsell opportunities")
        
        if not recommendations:
            recommendations = [
                "Account health is stable",
                "Continue regular engagement cadence",
                "Monitor for any changes in usage or support needs",
            ]
        
        pdf.set_font('Helvetica', '', 12)
        for rec in recommendations:
            pdf.cell(0, 8, f"• {rec}", ln=True)
        
        pdf.ln(20)
        pdf.set_font('Helvetica', 'I', 11)
        pdf.set_text_color(100, 116, 139)
        csm = self.account.get("csm_name", "your CSM")
        pdf.cell(0, 10, f"Questions? Contact {csm}", align='C')
